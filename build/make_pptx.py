# -*- coding: utf-8 -*-
"""
PM 역량인증 발표자료 생성 스크립트
기준 문서: 발표자료/발표자료_텍스트초안_v4.md (9슬라이드)
디자인 톤: 오피스/엔터프라이즈 보고형, 16:9
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
import copy

# ---------- 디자인 시스템 ----------
NAVY = RGBColor(0x1F, 0x2D, 0x50)        # 메인 헤더/타이틀
NAVY_DARK = RGBColor(0x13, 0x1C, 0x33)   # 가장 진한 강조
BLUE = RGBColor(0x2E, 0x5A, 0x9E)        # 포인트 블루
BLUE_LIGHT = RGBColor(0xEA, 0xF0, 0xF9)  # 연한 박스 배경
GRAY_DARK = RGBColor(0x3A, 0x3F, 0x47)   # 본문 텍스트
GRAY_MID = RGBColor(0x6B, 0x72, 0x7C)    # 보조 텍스트
GRAY_LIGHT = RGBColor(0xE6, 0xE8, 0xEB)  # 구분선/연한 배경
GRAY_BG = RGBColor(0xF5, 0xF6, 0xF8)     # 카드 배경
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)       # 긍정/달성
RED = RGBColor(0xB0, 0x2A, 0x37)         # 경고/한계
GOLD = RGBColor(0xB8, 0x8A, 0x2E)        # 강조 포인트(과한 사용 금지)

FONT = "맑은 고딕"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------- 헬퍼 ----------

def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, shadow=False, round_=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    if round_:
        try:
            shp.adjustments[0] = 0.04
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, color=GRAY_DARK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.0,
             wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return tb


def add_multirun_para(tf, runs, align=PP_ALIGN.LEFT, line_spacing=1.0, space_after=0, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(space_after)
    for txt, size, color, bold, italic in runs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = FONT
        r.font.color.rgb = color
    return p


def add_line(slide, x1, y1, x2, y2, color=GRAY_LIGHT, w=Pt(1)):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = w
    ln.shadow.inherit = False
    return ln


def add_arrow_down(slide, cx, y, size=Inches(0.22), color=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, cx - size/2, y, size, size)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_arrow_right(slide, x, cy, size=Inches(0.26), color=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, cy - size*0.35, size, size*0.7)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def style_table(table, header_fill=NAVY, header_color=WHITE, font_size=11,
                 body_font_size=11, align_center_cols=None, col_widths=None,
                 row_h=None):
    align_center_cols = align_center_cols or []
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r_idx, row in enumerate(table.rows):
        if row_h:
            row.height = row_h
        for c_idx, cell in enumerate(row.cells):
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c_idx in align_center_cols else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(header_font if False else (font_size if r_idx == 0 else body_font_size))
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                for p in tf.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for r in p.runs:
                        r.font.color.rgb = header_color
                        r.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else GRAY_BG
                for p in tf.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = GRAY_DARK


def fill_table(slide, x, y, w, h, data, col_widths=None, font_size=11,
               header_font_size=11, align_center_cols=None, row_h=None):
    rows = len(data)
    cols = len(data[0])
    gshape = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = gshape.table
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
    style_table(table, font_size=header_font_size, body_font_size=font_size,
                align_center_cols=align_center_cols, col_widths=col_widths, row_h=row_h)
    return table


def slide_header(slide, eyebrow, title, page_no, total=9):
    # 상단 네이비 바
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=NAVY)
    if eyebrow:
        add_text(slide, Inches(0.55), Inches(0.28), Inches(8), Inches(0.3),
                  eyebrow, size=12, color=BLUE, bold=True)
    add_text(slide, Inches(0.55), Inches(0.55) if eyebrow else Inches(0.4),
              Inches(11.5), Inches(0.6),
              title, size=24, color=NAVY, bold=True)
    add_line(slide, Inches(0.55), Inches(1.18), Inches(12.78), Inches(1.18),
              color=GRAY_LIGHT, w=Pt(1))
    # 페이지 번호 / 푸터
    add_text(slide, Inches(12.0), Inches(7.13), Inches(1.0), Inches(0.3),
              f"{page_no:02d} / {total:02d}", size=10, color=GRAY_MID, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.55), Inches(7.13), Inches(6), Inches(0.3),
              "PM 역량인증 발표 — On-Device AI Hub 스마트 플러그 PoC", size=9, color=GRAY_MID)


def causal_footer(slide, text):
    box = add_rect(slide, Inches(0.55), Inches(6.62), Inches(12.23), Inches(0.42),
                    fill=BLUE_LIGHT, round_=True)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(6.62), Inches(11.7), Inches(0.42))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "▶  " + text
    r.font.size = Pt(11.5)
    r.font.italic = True
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = NAVY


# ======================================================================
# S1. 시작 — 무엇을 보여줄 것인가
# ======================================================================
def build_s1():
    slide = add_slide()
    set_bg(slide, WHITE)
    # 좌측 네이비 패널
    add_rect(slide, 0, 0, Inches(4.6), SLIDE_H, fill=NAVY)
    add_rect(slide, Inches(4.6), 0, Inches(0.06), SLIDE_H, fill=GOLD)

    add_text(slide, Inches(0.6), Inches(0.7), Inches(3.6), Inches(0.4),
              "PM 역량 인증 발표", size=14, color=BLUE_LIGHT, bold=True)
    add_text(slide, Inches(0.6), Inches(1.25), Inches(3.7), Inches(2.0),
              "직함이 아니라,\n수행한 기능으로\n보여드립니다", size=30, color=WHITE, bold=True, line_spacing=1.15)

    add_line(slide, Inches(0.6), Inches(3.55), Inches(3.7), Inches(3.55), color=RGBColor(0x4A,0x57,0x78), w=Pt(1))

    add_text(slide, Inches(0.6), Inches(3.75), Inches(3.7), Inches(0.35),
              "On-Device AI Hub를 위한", size=14, color=RGBColor(0xC9,0xD2,0xE3))
    add_text(slide, Inches(0.6), Inches(4.05), Inches(3.7), Inches(0.4),
              "스마트 플러그 PoC 개발", size=17, color=WHITE, bold=True)

    add_text(slide, Inches(0.6), Inches(6.6), Inches(3.7), Inches(0.6),
              "황인휘  ·  총 9장  ·  15분", size=12, color=RGBColor(0xAE,0xB8,0xCC))

    # 우측 콘텐츠
    rx = Inches(5.1)
    rw = Inches(7.7)
    add_text(slide, rx, Inches(0.65), rw, Inches(0.4), "과제 개요", size=13, color=BLUE, bold=True)

    data = [
        ["항목", "내용"],
        ["과제 목표", "스마트 플러그 PoC SW 개발 및 On-Device AI Hub 연동 시연"],
        ["기간", "2025.04.01 ~ 2025.12.10 (계획 대비 2주 연장)"],
        ["팀 구성", "상근 4명(PL/PM 1 + SW개발 3), 비상근 3명 / 총 23MM"],
        ["관리 방식", "Waterfall / Confluence 기반"],
        ["핵심 제약", "공식 PM 타이틀 없이 PL이 관리 기능 수행"],
    ]
    fill_table(slide, rx, Inches(1.05), rw, Inches(2.5), data,
               col_widths=[Inches(1.7), Inches(6.0)], font_size=13, header_font_size=13)

    # 프레임 선언 박스
    box = add_rect(slide, rx, Inches(3.95), rw, Inches(1.65), fill=BLUE_LIGHT, round_=True)
    tb = slide.shapes.add_textbox(rx + Inches(0.35), Inches(4.12), rw - Inches(0.7), Inches(1.35))
    tf = tb.text_frame
    tf.word_wrap = True
    add_multirun_para(tf, [("이 발표는 직함이 아니라 관리 기능의 수행 여부로 PM 역량을 증빙합니다.",
                             15, NAVY, True, False)], line_spacing=1.3, first=True, space_after=8)
    add_multirun_para(tf, [("요구사항부터 종료까지, 한 단계의 결과가 다음 단계의 판단 근거가 되는 과정을 "
                             "그대로 보여드리겠습니다.", 14, GRAY_DARK, False, False)], line_spacing=1.3)

    add_text(slide, rx, Inches(5.85), rw, Inches(0.4), "발표 구조", size=13, color=BLUE, bold=True)
    steps = ["요구사항·추적", "계획·통합(CR)", "대시보드·예측", "참여·소통", "종료·교훈", "종합·Q&A"]
    n = len(steps)
    gap = Inches(0.12)
    box_w = Emu(int((rw - gap*(n-1)) / n))
    for i, s in enumerate(steps):
        bx = rx + i*(box_w + gap)
        add_rect(slide, bx, Inches(6.3), box_w, Inches(0.65), fill=GRAY_BG, line=GRAY_LIGHT, line_w=Pt(0.75), round_=True)
        tb2 = slide.shapes.add_textbox(bx + Inches(0.05), Inches(6.3), box_w - Inches(0.1), Inches(0.65))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf2.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = s
        r.font.size = Pt(9.5)
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = NAVY
        if i < n - 1:
            add_arrow_right(slide, bx + box_w + Emu(int(gap*0.15)), Inches(6.62), size=Inches(0.16), color=GRAY_MID)

    add_text(slide, Inches(0.55), Inches(7.13), Inches(8), Inches(0.3),
              "PM 역량인증 발표 — On-Device AI Hub 스마트 플러그 PoC", size=9, color=RGBColor(0x9A,0x9A,0x9A))
    add_text(slide, Inches(12.0), Inches(7.13), Inches(1.0), Inches(0.3),
              "01 / 09", size=10, color=GRAY_MID, align=PP_ALIGN.RIGHT)


# ======================================================================
# S2. 합의에서 시작하다 — 요구사항 정합성/추적성
# ======================================================================
def build_s2():
    slide = add_slide()
    set_bg(slide)
    slide_header(slide, "평가항목 ① 요구사항 정합성 · ② 추적성",
                 "44개의 합의, 그리고 그 합의를 끝까지 추적한 방법", 2)

    colL_x = Inches(0.55)
    colR_x = Inches(6.95)
    col_w = Inches(5.85)
    top_y = Inches(1.42)

    # 좌측: 정합성 funnel
    add_text(slide, colL_x, top_y, col_w, Inches(0.32), "정합성 — 무엇을 합의했는가", size=14, color=NAVY, bold=True)
    funnel = [
        ("Stakeholder Needs", ""),
        ("44개 REQ 확정", "PLG 27 + HUB 17"),
        ("우선순위 분류", "P0 / P1 / P2"),
        ("공식 승인", "TL(4/24) · iLab장(4/29)"),
    ]
    fy = top_y + Inches(0.45)
    fh = Inches(0.55)
    gap = Inches(0.18)
    for i, (main, sub) in enumerate(funnel):
        by = fy + i*(fh+gap)
        add_rect(slide, colL_x, by, col_w, fh, fill=(NAVY if i==1 else GRAY_BG),
                 line=GRAY_LIGHT, line_w=Pt(0.75), round_=True)
        tb = slide.shapes.add_textbox(colL_x+Inches(0.2), by, col_w-Inches(0.4), fh)
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = main
        r.font.size = Pt(13); r.font.bold = True; r.font.name = FONT
        r.font.color.rgb = WHITE if i==1 else NAVY
        if sub:
            r2 = p.add_run(); r2.text = "   " + sub
            r2.font.size = Pt(11); r2.font.name = FONT
            r2.font.color.rgb = (RGBColor(0xC9,0xD2,0xE3) if i==1 else GRAY_MID)
        if i < len(funnel)-1:
            add_arrow_down(slide, colL_x + col_w/2, by+fh+Emu(int(gap*0.15)), size=Inches(0.16))

    # 좌측 하단: 검증 기준
    vy = fy + 4*(fh+gap) + Inches(0.05)
    box = add_rect(slide, colL_x, vy, col_w, Inches(1.05), fill=BLUE_LIGHT, round_=True)
    tb = slide.shapes.add_textbox(colL_x+Inches(0.22), vy+Inches(0.08), col_w-Inches(0.44), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    add_multirun_para(tf, [("검증 기준 — 이유와 한계", 11.5, NAVY, True, False)], first=True, space_after=3)
    add_multirun_para(tf, [("PoC 범위 제외 기준(4.3 §1.3)에 따라 핵심 기능은 WBS 단위 DoD로 검증. "
                            "정식 테스트 설계(경계값·회귀 등)는 도입하지 못해 다음 단계 개선 과제로 남음.",
                            10.5, GRAY_DARK, False, False)], line_spacing=1.2)

    # 우측: 추적성 체인
    add_text(slide, colR_x, top_y, col_w, Inches(0.32), "추적성 — 합의를 어떻게 놓치지 않았는가", size=14, color=NAVY, bold=True)
    chain = [
        ("PLG_REQ_013", "기동어 인식 결과 처리 (P0)"),
        ("PLG_WBS_037 / 038 / 039", "구현"),
        ("DoD", "code commit + 동작 확인"),
        ("5.3 MS#4 기능 적합성 PASS", "TL 승인 (9/30)"),
    ]
    cy = top_y + Inches(0.45)
    ch = Inches(0.55)
    for i, (main, sub) in enumerate(chain):
        by = cy + i*(ch+gap)
        add_rect(slide, colR_x, by, col_w, ch, fill=(BLUE if i==3 else GRAY_BG),
                 line=GRAY_LIGHT, line_w=Pt(0.75), round_=True)
        tb = slide.shapes.add_textbox(colR_x+Inches(0.2), by, col_w-Inches(0.4), ch)
        tf = tb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = main
        r.font.size = Pt(12.5); r.font.bold = True; r.font.name = FONT
        r.font.color.rgb = WHITE if i==3 else NAVY
        r2 = p.add_run(); r2.text = "   " + sub
        r2.font.size = Pt(10.5); r2.font.name = FONT
        r2.font.color.rgb = (RGBColor(0xDC,0xE6,0xF5) if i==3 else GRAY_MID)
        if i < len(chain)-1:
            add_arrow_down(slide, colR_x + col_w/2, by+ch+Emu(int(gap*0.15)), size=Inches(0.16))

    ky = cy + 4*(ch+gap) + Inches(0.05)
    box2 = add_rect(slide, colR_x, ky, col_w, Inches(1.05), fill=GRAY_BG, line=GRAY_LIGHT, line_w=Pt(0.75), round_=True)
    tb = slide.shapes.add_textbox(colR_x+Inches(0.22), ky+Inches(0.08), col_w-Inches(0.44), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    add_multirun_para(tf, [("추적 요구사항 ", 13, GRAY_DARK, False, False), ("47개", 15, BLUE, True, False),
                           ("  ·  연결 WBS ", 13, GRAY_DARK, False, False), ("65개", 15, BLUE, True, False),
                           ("  ·  누락 ", 13, GRAY_DARK, False, False), ("0건", 15, GREEN, True, False)],
                          align=PP_ALIGN.CENTER, first=True, space_after=4)
    add_multirun_para(tf, [("CR-01/02 반영 후에도 전체 추적 유지 확인", 10.5, GRAY_MID, False, True)],
                          align=PP_ALIGN.CENTER)

    causal_footer(slide, "이 추적 체계가, 8월에 닥친 변경 요청의 영향을 즉시 계산할 수 있는 기반이 되었습니다.")


# ======================================================================
# S3. 8월, 계획에 없던 요청이 들어왔다 — 계획수립/통합 (CR-01/02)
# ======================================================================
def build_s3():
    slide = add_slide()
    set_bg(slide)
    slide_header(slide, "평가항목 ③ 계획수립 · ④ 통합",
                 "합의된 기준이 있었기에, 빠르게 따져볼 수 있었다", 3)

    top_y = Inches(1.42)
    add_text(slide, Inches(0.55), top_y, Inches(12.2), Inches(0.32),
              "8/11~12, iLab장·TL 요청으로 시연 시나리오 통합(CR-01) 및 플러그 4:1 확장(CR-02) 발생",
              size=12.5, color=GRAY_DARK, italic=True)

    add_text(slide, Inches(0.55), Inches(1.85), Inches(6), Inches(0.3),
              "영향도 분석", size=14, color=NAVY, bold=True)
    data = [
        ["영역", "영향 내용", "Score", "PM 1차 의견", "최종 결정"],
        ["범위·일정", "신규 구현 4주+, 검증 2주 추가", "25", "불수용", "수용"],
        ["성능", "멀티카메라 연동 네트워크 지연 리스크", "25", "불수용", "수용"],
        ["리스크", "타 과제 기개발 모듈 통합 일정 리스크", "25", "불수용", "수용"],
    ]
    fill_table(slide, Inches(0.55), Inches(2.2), Inches(7.6), Inches(1.35), data,
               col_widths=[Inches(1.2), Inches(3.4), Inches(0.75), Inches(1.15), Inches(1.1)],
               font_size=11.5, header_font_size=11.5, align_center_cols=[2,3,4])

    box = add_rect(slide, Inches(0.55), Inches(3.78), Inches(7.6), Inches(0.62), fill=BLUE_LIGHT, round_=True)
    tb = slide.shapes.add_textbox(Inches(0.78), Inches(3.78), Inches(7.2), Inches(0.62))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "결정 근거  "
    r.font.size = Pt(11.5); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = NAVY
    r2 = p.add_run(); r2.text = "조직 비즈니스 목표(사업부 기술 바잉 설득) 우선 판단에 따라 CCB에서 수용 결정 (8/12)"
    r2.font.size = Pt(11.5); r2.font.name = FONT; r2.font.color.rgb = GRAY_DARK

    # 수용 직후 4단계 조치
    add_text(slide, Inches(0.55), Inches(4.62), Inches(7.6), Inches(0.3),
              "수용 직후 4단계 조치", size=14, color=NAVY, bold=True)
    steps = [
        ("CCB 수용 결정", ""),
        ("영향도 분석 문서화", "6.1 §4"),
        ("기준선 재설정", "WBS 61→64개, 종료일 +14일"),
        ("추적매트릭스·대시보드 동기화", ""),
    ]
    sw = Inches(1.78)
    sgap = Inches(0.12)
    sy = Inches(5.0)
    for i,(main,sub) in enumerate(steps):
        sx = Inches(0.55) + i*(sw+sgap)
        add_rect(slide, sx, sy, sw, Inches(0.95), fill=(NAVY if i==2 else GRAY_BG),
                 line=GRAY_LIGHT, line_w=Pt(0.75), round_=True)
        tb = slide.shapes.add_textbox(sx+Inches(0.1), sy+Inches(0.08), sw-Inches(0.2), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = main
        r.font.size = Pt(11); r.font.bold = True; r.font.name = FONT
        r.font.color.rgb = WHITE if i==2 else NAVY
        if sub:
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run(); r2.text = sub
            r2.font.size = Pt(9); r2.font.name = FONT
            r2.font.color.rgb = RGBColor(0xC9,0xD2,0xE3) if i==2 else GRAY_MID
        if i < len(steps)-1:
            add_arrow_right(slide, sx+sw+Emu(int(sgap*0.1)), sy+Inches(0.475), size=Inches(0.2))

    # 우측: 일정 편차 (V1/V2)
    rx = Inches(8.5)
    add_text(slide, rx, Inches(1.85), Inches(4.28), Inches(0.3), "일정 편차 분석 (V1 → V2)", size=14, color=NAVY, bold=True)
    data2 = [
        ["시점", "계획", "실적", "편차"],
        ["V1 (8/1, CR 직전)", "WBS 12개", "15개", "+3개 선행"],
        ["V2 (12/9, 종료)", "WBS 64개 / 12/01", "61개 / 12/09", "-3개 / +8일"],
    ]
    fill_table(slide, rx, Inches(2.2), Inches(4.28), Inches(1.0), data2,
               col_widths=[Inches(1.4), Inches(1.15), Inches(1.05), Inches(0.68)],
               font_size=10.5, header_font_size=10.5, align_center_cols=[1,2,3])

    causal_footer(slide, "기준선이 바뀌자, 그 다음부터 우리가 추적해야 할 숫자 자체가 달라졌습니다.")


# ======================================================================
# S4. 숫자가 판단으로, 판단이 결과로 이어졌다 — 대시보드/미래예측
# ======================================================================
def build_s4():
    slide = add_slide()
    set_bg(slide)
    slide_header(slide, "평가항목 ⑤ 대시보드 현재 · ⑥ 미래예측",
                 "바뀐 기준 위에서, 무엇을 보고 무엇을 결정했는가", 4)

    top_y = Inches(1.42)
    add_text(slide, Inches(0.55), top_y, Inches(6), Inches(0.3), "KPI 3종 — 판단·조치 체인", size=14, color=NAVY, bold=True)

    chain1 = ["기능구현완료율 95%", "미달 5%는 시연 핵심 외 항목 → 시연 가능", "미완료 항목 시연 범위 외 분류", "시연 100% 완료"]
    chain2 = ["주요 결함 C급 2건 (A·B급 0건)", "허용 범위 내 Closed 처리", "A·B급 0건 유지"]

    def draw_chain(items, x, y, w, colors=None):
        h = Inches(0.5)
        gap = Inches(0.13)
        for i, txt in enumerate(items):
            by = y + i*(h+gap)
            fill = NAVY if i == len(items)-1 else GRAY_BG
            tcol = WHITE if i == len(items)-1 else NAVY
            add_rect(slide, x, by, w, h, fill=fill, line=GRAY_LIGHT, line_w=Pt(0.75), round_=True)
            tb = slide.shapes.add_textbox(x+Inches(0.15), by, w-Inches(0.3), h)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = txt
            r.font.size = Pt(10.5); r.font.bold = (i==0 or i==len(items)-1); r.font.name = FONT
            r.font.color.rgb = tcol
            if i < len(items)-1:
                add_arrow_down(slide, x+w/2, by+h+Emu(int(gap*0.1)), size=Inches(0.14))
        return y + len(items)*(h+gap)

    cw = Inches(2.85)
    draw_chain(chain1, Inches(0.55), top_y+Inches(0.4), cw)
    draw_chain(chain2, Inches(3.65), top_y+Inches(0.4), cw)

    # 일정 편차 표
    rx = Inches(6.85)
    add_text(slide, rx, top_y, Inches(5.95), Inches(0.3), "일정 편차 — V1/V2 비교 (CR 전후 연결)", size=14, color=NAVY, bold=True)
    data = [
        ["시점", "계획", "실적", "편차"],
        ["V1 (8/1, CR 발생 직전)", "WBS 12개", "15개", "+3개 선행"],
        ["V2 (12/9, 종료)", "WBS 64개 / 12/01", "61개 / 12/09", "-3개 / +8일"],
    ]
    fill_table(slide, rx, top_y+Inches(0.4), Inches(5.95), Inches(1.05), data,
               col_widths=[Inches(2.05), Inches(1.6), Inches(1.5), Inches(0.8)],
               font_size=11.5, header_font_size=11.5, align_center_cols=[1,2,3])

    # 인용 박스
    qy = top_y + Inches(2.05)
    box = add_rect(slide, rx, qy, Inches(5.95), Inches(1.7), fill=BLUE_LIGHT, round_=True)
    tb = slide.shapes.add_textbox(rx+Inches(0.28), qy+Inches(0.15), Inches(5.4), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    add_multirun_para(tf, [("7.5 대시보드 원문  ", 11, NAVY, True, False)], first=True, space_after=4)
    add_multirun_para(tf, [('"PLG 15개 완료(계획 대비 3개 선행)는 ', 11, GRAY_DARK, False, False),
                            ("CR-01/02 발생 후 대응 여력 확보에 기여함", 11, NAVY, True, True),
                            ('."', 11, GRAY_DARK, False, False)], line_spacing=1.25, space_after=6)
    add_multirun_para(tf, [("미래 예측 기법(EVM 등 잔여작업 기반 정량 예측)은 이번 과제에 도입하지 못한 부분이며, "
                            "S7 개선 기회에서 다룹니다.", 10.5, GRAY_MID, False, False)], line_spacing=1.2)

    causal_footer(slide, "이 KPI 판단이, 통제 가능한 영역에 자원을 쓰는 다음 선택으로 이어졌습니다.")


# ======================================================================
# S5. 막을 수 없는 것과, 흡수할 수 있는 것 — 정보배포/이해관계자참여
# ======================================================================
def build_s5():
    slide = add_slide()
    set_bg(slide)
    slide_header(slide, "평가항목 ⑦ 정보배포 · ⑧ 이해관계자참여",
                 "통제 가능한 영역에 자원을 쓰기로 했다", 5)

    top_y = Inches(1.42)
    add_text(slide, Inches(0.55), top_y, Inches(6), Inches(0.3), "참여 활동 실적", size=14, color=NAVY, bold=True)
    data = [
        ["채널", "대상", "주기", "실적"],
        ["주간 팀미팅", "TL, 개발", "매주", "31회 / 참석률 100%"],
        ["주간 랩보고", "iLab장, TL, 기획", "매주", "28회 / 참석률 100%"],
        ["게이트 리뷰", "전체 이해관계자", "최소 2회", "3회 진행"],
    ]
    fill_table(slide, Inches(0.55), top_y+Inches(0.4), Inches(6.0), Inches(1.5), data,
               col_widths=[Inches(1.4), Inches(1.8), Inches(1.0), Inches(1.8)],
               font_size=11, header_font_size=11, align_center_cols=[2,3])

    add_text(slide, Inches(0.55), top_y+Inches(2.1), Inches(6.0), Inches(0.5),
              "정보 배포: Collab(문서)·GitLab(소스) 중앙화 — 게이트 산출물 최종본 전체 반영(2.2)",
              size=10.5, color=GRAY_MID, italic=True)

    # 판단 흐름 (CR 이후)
    add_text(slide, Inches(0.55), top_y+Inches(2.65), Inches(6.0), Inches(0.3), "판단 — CR 이후", size=14, color=NAVY, bold=True)
    flow = [
        "CR-01/02는 정기 채널 외, 최상위 이해관계자의 권한 행사로 발생",
        "통제 가능 영역과 불가능 영역을 구분",
        "예방(채널·빈도 조정) → 효용 없음으로 판단, 폐기",
        "선택: S4에서 확보된 +3개 선행분을 버퍼로 흡수",
        "추가 인력 없이 기존 4인으로 대응 완료",
    ]
    fy = top_y + Inches(3.0)
    fh = Inches(0.46)
    gap = Inches(0.08)
    for i, txt in enumerate(flow):
        by = fy + i*(fh+gap)
        is_last = i == len(flow)-1
        add_rect(slide, Inches(0.55), by, Inches(6.0), fh, fill=(NAVY if is_last else GRAY_BG),
                 line=GRAY_LIGHT, line_w=Pt(0.75), round_=True)
        tb = slide.shapes.add_textbox(Inches(0.75), by, Inches(5.6), fh)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = txt
        r.font.size = Pt(10.5); r.font.name = FONT
        r.font.bold = is_last
        r.font.color.rgb = WHITE if is_last else GRAY_DARK
        if i < len(flow)-1:
            add_arrow_down(slide, Inches(0.55)+Inches(3.0), by+fh+Emu(int(gap*0.05)), size=Inches(0.13))

    # 우측: 통제 가능/불가능 시각 구분 카드
    rx = Inches(7.0)
    add_text(slide, rx, top_y, Inches(5.78), Inches(0.3), "통제 가능 / 불가능 영역 구분", size=14, color=NAVY, bold=True)

    card_w = Inches(2.8)
    card_h = Inches(2.0)
    cy = top_y + Inches(0.45)
    # 불가능 카드
    add_rect(slide, rx, cy, card_w, card_h, fill=GRAY_BG, line=GRAY_LIGHT, line_w=Pt(0.75), round_=True)
    tb = slide.shapes.add_textbox(rx+Inches(0.2), cy+Inches(0.18), card_w-Inches(0.4), card_h-Inches(0.36))
    tf = tb.text_frame; tf.word_wrap=True
    add_multirun_para(tf, [("통제 불가능", 13, RED, True, False)], first=True, space_after=8)
    add_multirun_para(tf, [("최상위 이해관계자의 직접 권한 행사(CR-01/02 발생 경로)", 11, GRAY_DARK, False, False)],
                          line_spacing=1.25, space_after=10)
    add_multirun_para(tf, [("✗ 채널·빈도 조정으로 예방", 10.5, GRAY_MID, False, False)], line_spacing=1.2, space_after=2)
    add_multirun_para(tf, [("→ 효용 없음으로 판단, 폐기", 10.5, GRAY_MID, False, True)], line_spacing=1.2)

    rx2 = rx + card_w + Inches(0.18)
    add_rect(slide, rx2, cy, card_w, card_h, fill=BLUE_LIGHT, line=BLUE, line_w=Pt(1), round_=True)
    tb = slide.shapes.add_textbox(rx2+Inches(0.2), cy+Inches(0.18), card_w-Inches(0.4), card_h-Inches(0.36))
    tf = tb.text_frame; tf.word_wrap=True
    add_multirun_para(tf, [("통제 가능", 13, GREEN, True, False)], first=True, space_after=8)
    add_multirun_para(tf, [("일정·자원 버퍼의 확보와 배분", 11, GRAY_DARK, False, False)],
                          line_spacing=1.25, space_after=10)
    add_multirun_para(tf, [("✔ +3개 선행분을 버퍼로 흡수", 10.5, NAVY, True, False)], line_spacing=1.2, space_after=2)
    add_multirun_para(tf, [("→ 추가 인력 없이 기존 4인 대응", 10.5, NAVY, False, True)], line_spacing=1.2)

    causal_footer(slide, "이 버퍼 흡수 판단이, 종료 시점에 \"무엇을 기준으로 끝났다고 할 것인가\"의 답이 됩니다.")


# ======================================================================
# S6. 끝났다고 말할 수 있는 근거 — 종료
# ======================================================================
def build_s6():
    slide = add_slide()
    set_bg(slide)
    slide_header(slide, "평가항목 ⑨ 종료", "무엇을 먼저 보고 종료를 선언했는가", 6)

    top_y = Inches(1.5)

    # 핵심 지표 강조 카드 4개 (상단 비주얼 요약)
    cards = [("핵심 시연 시나리오", "100%", "달성"), ("A·B급 결함", "0건", "목표 달성"),
             ("이해관계자 승인", "완료", "11/30"), ("WBS 완료율(보조)", "95.3%", "미완 3개")]
    cw = Inches(2.93)
    cgap = Inches(0.16)
    ch = Inches(1.15)
    for i,(label,val,sub) in enumerate(cards):
        cx = Inches(0.55) + i*(cw+cgap)
        is_aux = (i == 3)
        add_rect(slide, cx, top_y, cw, ch, fill=(GRAY_BG if is_aux else NAVY), round_=True)
        tb = slide.shapes.add_textbox(cx, top_y+Inches(0.12), cw, ch-Inches(0.24))
        tf = tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.size=Pt(26); r.font.bold=True; r.font.name=FONT
        r.font.color.rgb = GRAY_MID if is_aux else WHITE
        p2 = tf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = label
        r2.font.size=Pt(11); r2.font.bold=True; r2.font.name=FONT
        r2.font.color.rgb = GRAY_DARK if is_aux else RGBColor(0xE3,0xE8,0xF2)
        p3 = tf.add_paragraph(); p3.alignment=PP_ALIGN.CENTER
        r3 = p3.add_run(); r3.text = sub
        r3.font.size=Pt(9.5); r3.font.italic=True; r3.font.name=FONT
        r3.font.color.rgb = GRAY_MID if is_aux else RGBColor(0xAE,0xC0,0xDE)

    table_y = top_y + ch + Inches(0.35)
    add_text(slide, Inches(0.55), table_y, Inches(8), Inches(0.32),
              "종료 판단 — 핵심 지표 선행, 보조 지표 후행", size=14, color=NAVY, bold=True)

    data = [
        ["기준", "목표", "결과"],
        ["핵심 시연 시나리오 (CR로 확장된 범위)", "음성제어·인터컴·멀티플러그 4:1 통합", "100% 달성"],
        ["A·B급 결함", "0건", "0건"],
        ["이해관계자 최종 승인", "iLab장 승인", "완료 (11/30)"],
        ["WBS 완료율 (보조 지표)", "100%", "95.3% (미완료 3개, 고도화)"],
    ]
    table = fill_table(slide, Inches(0.55), table_y+Inches(0.42), Inches(12.23), Inches(2.3), data,
               col_widths=[Inches(4.3), Inches(4.43), Inches(3.5)],
               font_size=13, header_font_size=13, align_center_cols=[1,2])
    for r_idx in [1,2,3]:
        cell = table.cell(r_idx, 2)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.color.rgb = GREEN
    for c_idx in range(3):
        cell = table.cell(4, c_idx)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = GRAY_MID
                run.font.italic = True

    add_text(slide, Inches(0.55), table_y+Inches(3.05), Inches(8), Inches(0.3),
              "종료 기준서(1.4) §2.4 — iLab장 최종 승인 완료", size=11, color=GRAY_MID, italic=True)

    causal_footer(slide, "이 시연 범위 자체가 S3의 CR 결정으로 확장된 결과입니다. 변경이 곧 이 과제가 도달한 최종 목표의 정의였습니다.")


# ======================================================================
# S7. 거기서 무엇을 배웠는가 — 교훈/회고
# ======================================================================
def build_s7():
    slide = add_slide()
    set_bg(slide)
    slide_header(slide, "평가항목 ⑩ 교훈/회고", "같은 일이 또 생긴다면, 무엇을 먼저 바꿀 것인가", 7)

    top_y = Inches(1.45)
    add_text(slide, Inches(0.55), top_y, Inches(8), Inches(0.3), "핵심 교훈 3건", size=14, color=NAVY, bold=True)
    data = [
        ["ID", "영역", "핵심 내용"],
        ["LL_01", "방법론", "요구사항 변동성 높은 과제는 Agile/Hybrid가 적합"],
        ["LL_05", "일정", "CR 후 신규 WBS 세분화 원칙을 PM이 리더십으로 관철해야 함"],
        ["LL_11", "이해관계자", "최상위 이해관계자 변경 리스크는 전략 조정이 아닌 버퍼로 흡수해야 함"],
    ]
    fill_table(slide, Inches(0.55), top_y+Inches(0.4), Inches(12.23), Inches(1.55), data,
               col_widths=[Inches(1.3), Inches(2.0), Inches(8.93)],
               font_size=12.5, header_font_size=12.5, align_center_cols=[0,1])

    add_text(slide, Inches(0.55), top_y+Inches(2.25), Inches(8), Inches(0.3),
              "정직하게 인정하는 개선 기회", size=14, color=NAVY, bold=True)

    gy = top_y + Inches(2.65)
    gw = Inches(5.95)
    gh = Inches(2.85)
    cards = [
        ("잔여작업 기반 정량 예측(EVM 등) 미적용",
         "완료시점·품질 리스크 사전예측 체계 부재가 개선 과제",
         "⑥ 미래예측 — S8 Predict 카드에서 보완 설계"),
        ("자동 알림·의사소통관리계획 부재",
         "Collab/GitLab 중앙화에는 의존했으나 배포 체계는 미문서화",
         "⑦ 정보배포 — S8 Visibility 카드에서 보완 설계"),
    ]
    for i,(top,bottom,tag) in enumerate(cards):
        gx = Inches(0.55) + i*(gw+Inches(0.33))
        add_rect(slide, gx, gy, gw, gh, fill=GRAY_BG, line=GRAY_LIGHT, line_w=Pt(0.75), round_=True)
        add_rect(slide, gx, gy, gw, Inches(0.85), fill=BLUE_LIGHT, round_=True)
        tb = slide.shapes.add_textbox(gx+Inches(0.3), gy+Inches(0.12), gw-Inches(0.6), Inches(0.65))
        tf = tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = top
        r.font.size=Pt(13); r.font.bold=True; r.font.name=FONT; r.font.color.rgb=NAVY
        add_arrow_down(slide, gx+gw/2, gy+Inches(0.95), size=Inches(0.18))
        tb2 = slide.shapes.add_textbox(gx+Inches(0.3), gy+Inches(1.25), gw-Inches(0.6), Inches(0.85))
        tf2 = tb2.text_frame; tf2.word_wrap=True
        p2 = tf2.paragraphs[0]; p2.alignment=PP_ALIGN.LEFT; p2.line_spacing=1.25
        r2 = p2.add_run(); r2.text = bottom
        r2.font.size=Pt(12); r2.font.name=FONT; r2.font.color.rgb=GRAY_DARK
        add_line(slide, gx+Inches(0.3), gy+Inches(2.15), gx+gw-Inches(0.3), gy+Inches(2.15), color=GRAY_LIGHT, w=Pt(0.75))
        tb3 = slide.shapes.add_textbox(gx+Inches(0.3), gy+Inches(2.28), gw-Inches(0.6), Inches(0.5))
        tf3 = tb3.text_frame; tf3.word_wrap=True
        p3 = tf3.paragraphs[0]
        r3 = p3.add_run(); r3.text = tag
        r3.font.size=Pt(10.5); r3.font.italic=True; r3.font.bold=True; r3.font.name=FONT; r3.font.color.rgb=BLUE


# ======================================================================
# S8. PM 관점에서의 성과와 다음 설계
# ======================================================================
def build_s8():
    slide = add_slide()
    set_bg(slide)
    slide_header(slide, "전체 종합 — 성과 ①~⑩ / 개선 기회 ⑥⑦ 중심",
                 "증명된 것과, 다음에 설계할 것", 8)

    top_y = Inches(1.45)
    colw = Inches(5.95)
    # 좌측: 성과 표
    add_text(slide, Inches(0.55), top_y, colw, Inches(0.3), "성과 — 수치 근거", size=14, color=NAVY, bold=True)
    data = [
        ["관리 영역", "수행 수준", "근거"],
        ["요구사항 정의·승인", "44개, 공식 승인", "3.3"],
        ["변경 영향도 분석·통제", "정량 분석 + CCB", "6.1 §4"],
        ["리스크→이슈 전환", "7건 전수 추적, 전원 Closed", "9.1/9.4"],
        ["단계별 검증", "7차 마일스톤 전 단계 PASS", "5.3"],
    ]
    fill_table(slide, Inches(0.55), top_y+Inches(0.4), colw, Inches(2.1), data,
               col_widths=[Inches(2.1), Inches(2.85), Inches(1.0)],
               font_size=11.5, header_font_size=11.5, align_center_cols=[2])

    # 좌측 하단: 개선 기회 요약 (S7 연결)
    iy = top_y + Inches(2.75)
    box = add_rect(slide, Inches(0.55), iy, colw, Inches(1.55), fill=BLUE_LIGHT, round_=True)
    tb = slide.shapes.add_textbox(Inches(0.85), iy+Inches(0.18), colw-Inches(0.6), Inches(1.25))
    tf = tb.text_frame; tf.word_wrap=True
    add_multirun_para(tf, [("정직하게 인정한 개선 기회", 12.5, NAVY, True, False)], first=True, space_after=6)
    add_multirun_para(tf, [("⑥ 미래예측  ", 11, BLUE, True, False),
                            ("EVM 등 잔여작업 기반 정량 예측 미적용", 11, GRAY_DARK, False, False)],
                          line_spacing=1.3, space_after=3)
    add_multirun_para(tf, [("⑦ 정보배포  ", 11, BLUE, True, False),
                            ("자동 알림·의사소통관리계획 미문서화", 11, GRAY_DARK, False, False)],
                          line_spacing=1.3)

    # 우측: Predict / Visibility / Resilience 카드
    rx = Inches(6.85)
    add_text(slide, rx, top_y, Inches(5.93), Inches(0.3), "다음 설계 — Predict / Visibility / Resilience", size=14, color=NAVY, bold=True)
    cards = [
        ("Predict", "KPI 임계치·판단 규칙을 착수 시점에 명문화. 잔여 WBS 기반 완료시점·품질 리스크 정량 예측 도입"),
        ("Visibility", "REQ-WBS-검증 상태를 단일 화면에서 확인 가능하게 구성. 의사소통관리계획을 별도 수립"),
        ("Resilience", "최상위 이해관계자 변경 가능성을 착수 시점 리스크로 등록, 일정·자원에 10~15% 버퍼 명시적 반영"),
    ]
    ch = Inches(1.15)
    cgap = Inches(0.18)
    for i,(label,desc) in enumerate(cards):
        cy = top_y+Inches(0.42)+i*(ch+cgap)
        add_rect(slide, rx, cy, Inches(5.93), ch, fill=GRAY_BG, line=GRAY_LIGHT, line_w=Pt(0.75), round_=True)
        add_rect(slide, rx, cy, Inches(0.08), ch, fill=BLUE)
        tb = slide.shapes.add_textbox(rx+Inches(0.3), cy+Inches(0.1), Inches(1.3), ch-Inches(0.2))
        tf = tb.text_frame; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = label
        r.font.size=Pt(15); r.font.bold=True; r.font.name=FONT; r.font.color.rgb=NAVY
        tb2 = slide.shapes.add_textbox(rx+Inches(1.65), cy+Inches(0.1), Inches(4.1), ch-Inches(0.2))
        tf2 = tb2.text_frame; tf2.word_wrap=True; tf2.vertical_anchor=MSO_ANCHOR.MIDDLE
        p2 = tf2.paragraphs[0]; p2.line_spacing=1.25
        r2 = p2.add_run(); r2.text = desc
        r2.font.size=Pt(11); r2.font.name=FONT; r2.font.color.rgb=GRAY_DARK

    causal_footer(slide, "이번 과제에서는 우연히 확보된 선행분이 버퍼 역할을 했습니다. 다음에는 그 버퍼를 의도적으로 설계하겠습니다.")


# ======================================================================
# S9. 정리 + Q&A 브리지
# ======================================================================
def build_s9():
    slide = add_slide()
    set_bg(slide, NAVY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=GOLD)

    add_text(slide, Inches(0.55), Inches(0.55), Inches(11), Inches(0.6), "정리", size=26, color=WHITE, bold=True)

    box = add_rect(slide, Inches(0.55), Inches(1.35), Inches(12.23), Inches(1.4),
                    fill=RGBColor(0x29,0x37,0x5C), round_=True)
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap=True
    add_multirun_para(tf, [("요구사항에서 종료까지, 한 단계의 결과가 다음 단계의 판단 근거가 되는 흐름을 수치로 증빙했습니다.",
                            15, WHITE, True, False)], line_spacing=1.3, first=True, space_after=6)
    add_multirun_para(tf, [("부족한 부분은 정확히 식별했고, 다음 설계까지 구체화했습니다.",
                            14, RGBColor(0xC9,0xD2,0xE3), False, False)], line_spacing=1.3)

    # 3개 강조 포인트
    points = [
        ("①", "자기 인식", "약점을 회피하지 않고 직시"),
        ("②", "해결 능력", "보완 방향이 구체적이고 실행 가능"),
        ("③", "균형 역량", "장단점이 고르게 분포된 PM 프로파일"),
    ]
    pw = Inches(3.95)
    py = Inches(3.05)
    for i,(num,title,desc) in enumerate(points):
        px = Inches(0.55) + i*(pw+Inches(0.19))
        add_rect(slide, px, py, pw, Inches(1.5), fill=RGBColor(0x29,0x37,0x5C), line=RGBColor(0x4A,0x57,0x78), line_w=Pt(0.75), round_=True)
        tb = slide.shapes.add_textbox(px+Inches(0.25), py+Inches(0.15), pw-Inches(0.5), Inches(1.2))
        tf = tb.text_frame; tf.word_wrap=True
        add_multirun_para(tf, [(num+"  ", 16, GOLD, True, False), (title, 15, WHITE, True, False)],
                              first=True, space_after=6)
        add_multirun_para(tf, [(desc, 11.5, RGBColor(0xC9,0xD2,0xE3), False, False)], line_spacing=1.25)

    # Q&A 브리지
    add_text(slide, Inches(0.55), Inches(4.85), Inches(8), Inches(0.32),
              "Q&A — 질문 주시면 아래 순서로 근거를 바로 보여드리겠습니다", size=14, color=GOLD, bold=True)

    qa = [
        ("1", "CR 판단·영향도 분석 근거", "6.1 / 7.5"),
        ("2", "REQ → WBS → 검증 추적 체인", "4.2 / 4.3 / 5.3"),
        ("3", "KPI 판단 근거 및 종료 승인", "1.5 / 1.4"),
    ]
    qw = Inches(3.95)
    qy = Inches(5.3)
    for i,(num,title,ref) in enumerate(qa):
        qx = Inches(0.55) + i*(qw+Inches(0.19))
        add_rect(slide, qx, qy, qw, Inches(1.1), fill=WHITE, round_=True)
        tb = slide.shapes.add_textbox(qx+Inches(0.22), qy+Inches(0.12), qw-Inches(0.44), Inches(0.9))
        tf = tb.text_frame; tf.word_wrap=True
        add_multirun_para(tf, [(num+". ", 13, BLUE, True, False), (title, 12.5, NAVY, True, False)],
                              line_spacing=1.2, first=True, space_after=4)
        add_multirun_para(tf, [(ref, 11, GRAY_MID, False, False)])

    add_text(slide, Inches(12.0), Inches(7.13), Inches(1.0), Inches(0.3),
              "09 / 09", size=10, color=RGBColor(0x9A,0xA4,0xBA), align=PP_ALIGN.RIGHT)


build_s1()
build_s2()
build_s3()
build_s4()
build_s5()
build_s6()
build_s7()
build_s8()
build_s9()

OUT_PATH = "../PM역량인증_황인휘_발표자료.pptx"
prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")

